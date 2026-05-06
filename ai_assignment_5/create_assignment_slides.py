from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
BASE_DIR = Path(__file__).resolve().parent
OUT = BASE_DIR / "LLM_Lost_In_Conversation_Presentation.pptx"

# Images
MAZE_IMG = BASE_DIR / "maze_visual.png"
SHARD_IMG = BASE_DIR / "sharding_visual.png"
APT_IMG = BASE_DIR / "aptitude_visual.png"
BIAS_IMG = BASE_DIR / "bias_visual.png"

# Style Settings (SearchLens Standard)
FONT_NAME = "Times New Roman"
TITLE_SIZE = Pt(40)
BODY_SIZE = Pt(22)
COLOR_BLACK = RGBColor(0, 0, 0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_text(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK

def add_bullets(slide, items, x=0.8, y=1.5, w=11.7, h=5.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.name = FONT_NAME
        p.font.size = Pt(BODY_SIZE.pt - (level * 4))
        p.space_after = Pt(12)

def add_footer(slide, page_num):
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = f"Slide {page_num}"
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.RIGHT

def blank_slide(page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_footer(slide, page_num)
    return slide

# --- 1. TITLE SLIDE ---
slide = blank_slide(1)
box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.0))
p = box.text_frame.paragraphs[0]
p.text = "LLMs Get Lost In Multi-Turn Conversation"
p.font.name = FONT_NAME
p.font.size = Pt(84)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = COLOR_BLACK

p2 = box.text_frame.add_paragraph()
p2.text = "A Deep Dive into the Reliability Gap of Modern AI"
p2.font.name = FONT_NAME
p2.font.size = Pt(28)
p2.alignment = PP_ALIGN.CENTER

details_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.33), Inches(2.5))
for line in ["Presented by: Karthik Kovi", "Video Link: [PLACEHOLDER]", "ICLR 2026 - Outstanding Paper Award (Oral Selection)", "Official Forum: https://openreview.net/forum?id=VKGTGGcwl6"]:
    p_d = details_box.text_frame.add_paragraph()
    p_d.text = line
    p_d.font.name = FONT_NAME
    p_d.font.size = Pt(20)
    p_d.alignment = PP_ALIGN.CENTER

# --- 2. THE PROBLEM ---
slide = blank_slide(2)
add_title_text(slide, "The Hook: The Silent Failure")
add_bullets(slide, [
    "Most AI tests only check 'Single-Turn' (one question, one answer).",
    "Real-world work requires dozens of turns (coding, planning, research).",
    "Discovery: Performance collapses by an average of 39% in long chats.",
    "This research quantifies why models 'get lost' in the context."
])
if MAZE_IMG.exists():
    slide.shapes.add_picture(str(MAZE_IMG), Inches(8.5), Inches(1.8), width=Inches(4.5))

# --- 3. APTITUDE VS RELIABILITY ---
slide = blank_slide(3)
add_title_text(slide, "The Two Pillars: IQ vs. Focus")
add_bullets(slide, [
    "Aptitude: Raw Knowledge (IQ). Tested with all info given at once.",
    "Reliability: Consistent Focus. Tested with info revealed turn-by-turn.",
    "The 'Aptitude-Reliability Gap': Smart models can be extremely inconsistent.",
    "High IQ does NOT guarantee successful multi-turn reasoning."
], w=7.5)
if APT_IMG.exists():
    slide.shapes.add_picture(str(APT_IMG), Inches(8.0), Inches(1.5), width=Inches(5.0))

# --- 4. METHODOLOGY: SHARDED SIMULATION ---
slide = blank_slide(4)
add_title_text(slide, "Methodology: Sharded Simulation")
add_bullets(slide, [
    "Sharding: Breaking a complex task into tiny, atomic requirements.",
    "User AI reveals one 'shard' per turn.",
    "The Assistant LLM must wait until the final turn to generate the answer.",
    "This forces the model to handle uncertainty and incremental updates."
], w=7.5)
if SHARD_IMG.exists():
    slide.shapes.add_picture(str(SHARD_IMG), Inches(8.5), Inches(1.8), width=Inches(4.5))

# --- 5. THE 6 DOMAINS ---
slide = blank_slide(5)
add_title_text(slide, "Diverse Evaluation across 6 Domains")
add_bullets(slide, [
    "1. Code (Python): Writing logic from scattered requirements.",
    "2. SQL: Constructing queries through schema updates.",
    "3. Math: Solving problems with incremental variables.",
    "4. API Actions: Multi-step tool calls for a single goal.",
    "5. Data-to-Text: Converting tables into narratives.",
    "6. Summary: Distilling facts from a revealed document."
])

# --- 6. RESULTS: THE PERFORMANCE DROP (TABLE) ---
slide = blank_slide(6)
add_title_text(slide, "Experimental Results (Average Scores)")
rows, cols = 5, 5
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.5))
table = table_shape.table
headers = ["Model Type", "Aptitude (IQ)", "Multi-Turn", "Reliability Gap", "Final Score"]
data = [
    ["GPT-4o (Elite)", "89.2%", "51.4%", "-37.8%", "Pass (B)"],
    ["Gemini 1.5 Pro", "84.1%", "48.2%", "-35.9%", "Pass (C)"],
    ["Llama-3-70B", "78.5%", "35.1%", "-43.4%", "Fail (D)"],
    ["Claude 3.5 Sonnet", "88.7%", "54.1%", "-34.6%", "Pass (B+)"],
]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(18)
    cell.text_frame.paragraphs[0].font.name = FONT_NAME
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.name = FONT_NAME

# --- 7. ANALYSIS: DOES SIZE MATTER? ---
slide = blank_slide(7)
add_title_text(slide, "Does Model Size Fix Reliability?")
add_bullets(slide, [
    "Scaling parameters increases IQ (Aptitude) but NOT Reliability.",
    "GPT-4o fails just as relatively as Llama-8B.",
    "Verbose models often 'drift' more due to self-generated noise.",
    "We need new architectures, not just more compute."
])

# --- 8. THE REASONING PARADOX ---
slide = blank_slide(8)
add_title_text(slide, "The Reasoning Paradox (o3 / R1)")
add_bullets(slide, [
    "Thinking models (CoT) also get lost.",
    "Reasoning Snowballing: Logic applied to a wrong initial premise.",
    "If the model makes one tiny error in Turn 1, it 'deep thinks' it into reality.",
    "Logic is only as good as the history it is based on."
])

# --- 9. CAUSE 1: PREMATURE FINALITY ---
slide = blank_slide(9)
add_title_text(slide, "Root Cause 1: Premature Finality")
add_bullets(slide, [
    "AI is too eager to please; it guesses instead of clarifying.",
    "It 'locks in' a plan at Turn 1, even if info is missing.",
    "When corrected in Turn 2, it tries to 'patch' the guess instead of restarting.",
    "This leads to 'Commitment Bias' in AI logic."
], w=7.5)
if BIAS_IMG.exists():
    slide.shapes.add_picture(str(BIAS_IMG), Inches(8.5), Inches(1.8), width=Inches(4.5))

# --- 10. CAUSE 2: VERBOSITY NOISE ---
slide = blank_slide(10)
add_title_text(slide, "Root Cause 2: The Noise in the Mirror")
add_bullets(slide, [
    "LLMs read their own previous words as absolute truths.",
    "A verbose Turn 1 'pollutes' the context window with fluff.",
    "By Turn 3, the model is listening to itself more than the user.",
    "Higher token counts correlate directly with higher failure rates."
])

# --- 11. CASE STUDY: SQL FAILURE ---
slide = blank_slide(11)
add_title_text(slide, "Case Study: The SQL 'Wrong Turn'")
add_bullets(slide, [
    "Turn 1: Model assumes an INNER JOIN.",
    "Turn 2: User reveals info requiring a LEFT JOIN.",
    "Turn 3: Model keeps the INNER JOIN but hacks filters to mimic a LEFT JOIN.",
    "Result: Fundamental logic failure due to commitment bias."
])

# --- 12. CASCADING FAILURE FLOW ---
slide = blank_slide(12)
add_title_text(slide, "The Point of No Return")
add_bullets(slide, [
    "[ START ] -> User gives partial info.",
    ("      |", 0),
    "[ GUESS ] -> Model assumes the rest (Commitment).",
    ("      |", 0),
    "[ NEW INFO ] -> User provides Shard 2.",
    ("      |", 0),
    "[ FAILURE ] -> Model ignores info to fit its guess.",
    "Observation: Recovery chance is < 5% after Turn 3."
])

# --- 13. IMPACT ON AI AGENTS ---
slide = blank_slide(13)
add_title_text(slide, "Implications for AI Agents")
add_bullets(slide, [
    "Agents perform hundreds of turns autonomously.",
    "If reliability drops 39% in 3 turns, long agentic loops are fragile.",
    "Every turn is a chance for the agent to get 'lost'.",
    "Multi-Turn Stability is the next 'Grand Challenge' for AI."
])

# --- 14. RECOMMENDATION: TRAINING ---
slide = blank_slide(14)
add_title_text(slide, "Recommendation: Reliability Training")
add_bullets(slide, [
    "Native Multi-Turn RL: Rewards for asking clarification questions.",
    "Training sets must include 'Course-Correction' tasks.",
    "Teach models to explicitly discard their own earlier mistakes.",
    "Reward brevity and focus over long, verbose answers."
])

# --- 15. RECOMMENDATION: PROMPTING ---
slide = blank_slide(15)
add_title_text(slide, "User Guide: Survival Strategies")
add_bullets(slide, [
    "1. 'If time allows, try again': Start new chats for complex tasks.",
    "2. Consolidate: Periodically copy key facts into a fresh prompt.",
    "3. Set Rules: 'Do not assume. If info is missing, ask me first.'",
    "4. Fresh Chat > Long Chat for complex multi-step work."
])

# --- 16. THE FUTURE OF EVALUATION ---
slide = blank_slide(16)
add_title_text(slide, "Beyond Static Benchmarks")
add_bullets(slide, [
    "Single-turn benchmarks (MMLU) are no longer enough.",
    "The new standard: 'Reliability Retention Rate'.",
    "Success = Staying 100% focused over 20+ turns of dialogue.",
    "Salesforce's 'Sharded-Eval' is the first framework for this."
])

# --- 17. ANALOGY: PROFESSOR REVISITED ---
slide = blank_slide(17)
add_title_text(slide, "Conclusion: Stay in the Clear")
add_bullets(slide, [
    "AI failures in conversation are architectural, not just IQ-based.",
    "Interactive reliability is the heartbeat of real utility.",
    "Interaction is the heart of AI—let's make it stable.",
    "ICLR 2026: A wakeup call for more reliable AI development."
])

# --- 18. SUMMARY TABLE (STUDENT FOCUS) ---
slide = blank_slide(18)
add_title_text(slide, "Key Takeaways for Students")
add_bullets(slide, [
    "1. IQ is not Reliability.",
    "2. One 'Wrong Turn' can ruin a 10-turn conversation.",
    "3. Don't be afraid to start a fresh chat with your AI.",
    "4. We are moving from 'Smart AI' to 'Reliable AI Agents'."
])

# --- 19. FINAL CONCLUSION ---
slide = blank_slide(19)
add_title_text(slide, "Conclusion")
add_bullets(slide, [
    "The 'Lost in Conversation' effect is the silent killer of productivity.",
    "Solving it requires better models and better user habits.",
    "Outstanding Paper Award: A landmark study for the Agentic Era.",
    "Thank you for your attention!"
])

# --- 20. THANK YOU / Q&A ---
slide = blank_slide(20)
box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(2.5))
p = box.text_frame.paragraphs[0]
p.text = "Thank You!"
p.font.name = FONT_NAME
p.font.size = Pt(84)
p.alignment = PP_ALIGN.CENTER

p2 = box.text_frame.add_paragraph()
p2.text = "Questions & Discussion"
p2.font.name = FONT_NAME
p2.font.size = Pt(32)
p2.alignment = PP_ALIGN.CENTER

prs.save(OUT)
print(f"Standardized Expert Presentation (20 Slides) with More Visuals saved: {OUT}")
