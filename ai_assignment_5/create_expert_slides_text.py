from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

OUT = Path("/Users/karthikkovi/Documents/CSULB/AI/Assignment-5/Expert_Presentation_20_Slides.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(title, points, takeaway=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Points
    box2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.0))
    tf = box2.text_frame
    tf.word_wrap = True
    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = point
        p.font.size = Pt(22)
        p.space_after = Pt(10)
        if point.startswith("- "):
            p.level = 0
            p.text = point[2:]
        elif point.startswith("  - "):
            p.level = 1
            p.text = point[4:]
            
    if takeaway:
        box3 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7))
        p = box3.text_frame.paragraphs[0]
        p.text = f"💡 Insight: {takeaway}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 102, 0)
        p.alignment = PP_ALIGN.CENTER

# 20 Slides Content
slides_content = [
    ("LLMs Get Lost In Multi-Turn Conversation", ["ICLR 2026 - Outstanding Paper Award", "Presented by: Karthik Kovi", "[VIDEO LINK PLACEHOLDER]"], "The most prestigious research on AI reliability."),
    ("The Hook: The Silent Failure", ["- We trust AI for long chats, but evaluations only test single turns.", "- Discovery: 39% performance drop as conversations progress.", "- It's like a 'wrong turn' in a maze that the AI can't fix."], "IQ is high, but reliability is fragile."),
    ("Analogy: The Professor in the Fog", ["- Every turn of chat adds 'fog' (complexity).", "- The Professor starts guessing what's inside the fog.", "- Eventually, the AI is solving a problem it imagined, not yours."], "Early guesses lead to total confusion."),
    ("Aptitude vs. Reliability", ["- Aptitude: The 'Knowledge' (Raw IQ).", "- Reliability: The 'Focus' (Consistency).", "- Finding: Models have the IQ but lose the Focus."], "IQ alone doesn't win conversations."),
    ("Methodology: Sharded Simulation", ["- 'Sharding' breaks complex tasks into small pieces.", "- Pieces are revealed one-by-one.", "- This forces the AI to handle uncertainty and updates."], "Real-world interaction is never one-and-done."),
    ("System Architecture", ["- User Simulator reveals shards gradually.", "- Assistant LLM tries to solve the task.", "- System Evaluator verifies the final result."], "Automated testing across 200,000 interactions."),
    ("The 6 Generation Tasks", ["- Code (Python)", "- Database (SQL)", "- API Actions", "- Math", "- Data-to-Text", "- Summarization"], "The failure is universal across all domains."),
    ("The Results: 39% Collapse", ["- Universal drop in success rates.", "- Failure spikes as early as Turn 2.", "- Models become 'overconfident' in wrong assumptions."], "Benchmarks don't show the real-world truth."),
    ("Model Size Analysis", ["- Myth: 'Bigger models are more reliable.'", "- Reality: GPT-4o fails as often as Llama-8B.", "- Scaling IQ doesn't scale Reliability."], "Parameter count isn't a cure for focus loss."),
    ("The Reasoning Paradox", ["- 'Thinking' models (o3, R1) also get lost.", "- 'Reasoning Snowballing': Logic applied to wrong premises.", "- Thinking longer often reinforces earlier errors."], "Thought is only useful if the facts are right."),
    ("Root Cause: Premature Finality", ["- AI tries to solve the task too early.", "- It guesses instead of asking questions.", "- It 'locks in' a plan and refuses to change it."], "Eagerness to answer leads to failure."),
    ("Root Cause: Verbosity Noise", ["- AI reads its own chatter as ground truth.", "- Verbose answers fill the context with fluff.", "- Model forgets the user's instructions."], "Brevity is the soul of reliability."),
    ("Case Study: SQL Failure", ["- Turn 1: Model picks a JOIN type.", "- Turn 2: New info requires a change.", "- Turn 3: Model tries to 'hack' the old join instead of fixing it."], "Commitment bias is real in LLMs."),
    ("The Cascading Error State", ["- One wrong turn leads to another.", "- The 'Point of No Return' is usually Turn 2 or 3.", "- Recovery chance is near zero after a guess."], "Early accuracy is the only way to win."),
    ("Danger for AI Agents", ["- Agents do hundreds of turns.", "- If 3 turns drop success by 39%, 50 turns are very risky.", "- We must solve this for the 'Agentic Future'."], "Agent stability is the next big hurdle."),
    ("Expert Suggestion: Multi-Turn RL", ["- Penalize guessing in training.", "- Reward 'Clarification Questions'.", "- Train models to 'Change their Mind' when info updates."], "Architecture must favor reliability over IQ."),
    ("User Guide: Tips to Win", ["- 'If time allows, try again' (New Chat).", "- Consolidate info every few turns.", "- Tell the AI: 'Do not assume. Ask me.'"], "A fresh map beats a lost guide."),
    ("The Future: Reliability Evals", ["- Moving beyond simple IQ tests.", "- New metric: 'Reliability Retention Rate'.", "- Success = staying focused over 10+ turns."], "Reliability is the new Gold Standard."),
    ("Conclusion", ["- LLMs get lost because they guess and over-talk.", "- 39% drop is a major barrier to real utility.", "- Solution: Reliability training and better usage."], "Let's build AI that stays on track."),
    ("Q&A", ["Paper: 'LLMs Get Lost In Multi-Turn Conversation'", "ICLR 2026 Outstanding Paper", "Thank you for your time!"], "Any questions?")
]

for title, points, takeaway in slides_content:
    add_slide(title, points, takeaway)

prs.save(OUT)
print(f"Expert 20-slide text-only PPTX saved: {OUT}")
