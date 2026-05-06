from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
prs.save("simple_test.pptx")
print("Done")
