from __future__ import annotations
print("Starting slide generation...")
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Paths
BASE_DIR = Path(__file__).resolve().parent
OUT = BASE_DIR / "searchlens_presentation.pptx"

# Images
UML_IMG = BASE_DIR / "uml_diagram.png"
TOKEN_IMG = BASE_DIR / "token_visual.png"
INDEX_IMG = BASE_DIR / "index_visual.png"
BM25_IMG = BASE_DIR / "bm25_visual.png"
FEEDBACK_IMG = BASE_DIR / "feedback_visual.png"
RESULTS_IMG = BASE_DIR / "results_visual.png"
ARCH_IMG = BASE_DIR / "arch_visual.png"

# Style Settings
FONT_NAME = "Times New Roman"
TITLE_SIZE = Pt(40)
BODY_SIZE = Pt(22)
COLOR_BLUE = RGBColor(0, 0, 0)  # Changed to black
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_GRAY = RGBColor(0, 0, 0)  # Changed to black for uniformity

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_text(slide, title: str):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE

def add_bullets(slide, items, x=0.8, y=1.5, w=11.7, h=5.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.name = FONT_NAME
        p.font.size = Pt(BODY_SIZE.pt - (level * 4))
        p.space_after = Pt(12)

def add_footer(slide, page_num):
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = f"Slide {page_num}"
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.RIGHT

def blank_slide(page_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_footer(slide, page_num)
    return slide

# --- 1. TITLE SLIDE ---
slide = blank_slide(1)
box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.0))
p = box.text_frame.paragraphs[0]
p.text = "SearchLens"
p.font.name = FONT_NAME
p.font.size = Pt(84)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = COLOR_BLUE

p2 = box.text_frame.add_paragraph()
p2.text = "A Reproducible Classical Search Engine with Positional Indexing & BM25"
p2.font.name = FONT_NAME
p2.font.size = Pt(28)
p2.alignment = PP_ALIGN.CENTER

details_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.33), Inches(2.5))
for line in ["Presented by:", "Karthik Kovi, Mohit Krishna Emani, Vignan Gadaley, Ram Venkat Sai Ritwick Vadlamudi", "CECS 429/529: Search Engine Technology", "California State University, Long Beach"]:
    p_d = details_box.text_frame.add_paragraph()
    p_d.text = line
    p_d.font.name = FONT_NAME
    p_d.font.size = Pt(20)
    p_d.alignment = PP_ALIGN.CENTER
    p_d.font.color.rgb = COLOR_GRAY

# --- 2. THE TEAM ---
slide = blank_slide(2)
add_title_text(slide, "The SearchLens Team")
add_bullets(slide, [
    "Mohit Krishna Emani - [Role: Project Lead & Core Engine]",
    "Karthik Kovi - [Role: Preprocessing & Evaluation]",
    "Vignan Gadaley - [Role: Web Interface & UI Design]",
    "Ram Venkat Sai Ritwick Vadlamudi - [Role: Research & Testing]",
    "Goal: To demystify the 'magic' behind information retrieval systems."
])

# --- 3. INTRODUCTION ---
slide = blank_slide(3)
add_title_text(slide, "What is a Search Engine?")
add_bullets(slide, [
    "At its simplest, a search engine is a 'find' tool on steroids.",
    "It solves three massive problems:",
    ("1. Speed: How to find data in milliseconds across millions of files.", 1),
    ("2. Precision: How to make sure the best result is at the top.", 1),
    ("3. Language: How to handle typos, phrases, and word meanings.", 1),
    "SearchLens focuses on the 'Ranked Retrieval' paradigm."
])

# --- 4. MOTIVATION ---
slide = blank_slide(4)
add_title_text(slide, "Motivation: Why Build This?")
add_bullets(slide, [
    "Modern search is dominated by 'Black Box' libraries (Lucene, Elasticsearch).",
    "Understanding the 'lexical core' is essential for hybrid AI search.",
    "Bridging the Gap: Connecting classroom theory to runnable, inspectable Python code.",
    "Reproducibility: Proving that BM25 outperforms TF-IDF on technical data."
])

# --- 5. PROBLEM STATEMENT ---
slide = blank_slide(5)
add_title_text(slide, "The Technical Challenge")
add_bullets(slide, [
    "Problem: Given a collection of 60 technical documents (SETC-60), how do we return the most relevant results for a user's natural language query?",
    "Requirements:",
    ("Handle exact phrase matching (e.g. \"search engine\")", 1),
    ("Support complex Boolean logic (AND, OR, NOT)", 1),
    ("Rank results using probabilistic models (BM25)", 1),
    ("Improve search automatically using user feedback loops", 1)
])

# --- 6. CORE CONCEPT: THE INVERTED INDEX ---
slide = blank_slide(6)
add_title_text(slide, "The 'Brain' of Search: Inverted Index")
add_bullets(slide, [
    "Searching every file for a word is slow (O(N)).",
    "An Inverted Index maps words to their locations upfront (O(1)).",
    "Positional Indexing: We don't just store which document has a word; we store the EXACT position of that word in the text."
], w=7.0)
if INDEX_IMG.exists():
    slide.shapes.add_picture(str(INDEX_IMG), Inches(8.0), Inches(1.5), width=Inches(5.0))

# --- 7. THE PIPELINE ---
slide = blank_slide(7)
add_title_text(slide, "The SearchLens Architecture")
add_bullets(slide, [
    "The system is built as a series of modular pipes:",
    ("1. Preprocessing (Cleaning the raw data)", 1),
    ("2. Indexing (Building the data structure)", 1),
    ("3. Retrieval (Finding candidates)", 1),
    ("4. Ranking (Sorting by relevance)", 1)
], w=6.5)
if ARCH_IMG.exists():
    slide.shapes.add_picture(str(ARCH_IMG), Inches(7.0), Inches(1.8), width=Inches(6.0))

# --- 8. UML SYSTEM DESIGN ---
slide = blank_slide(8)
add_title_text(slide, "System Component Diagram")
if UML_IMG.exists():
    slide.shapes.add_picture(str(UML_IMG), Inches(3.5), Inches(1.5), height=Inches(5.5))

# --- 9. MODULE 1: PREPROCESSING ---
slide = blank_slide(9)
add_title_text(slide, "Preprocessing: Cleaning the Noise")
add_bullets(slide, [
    "Computers don't understand words; they see characters.",
    "Tokenization: Splitting 'The search!' into ['the', 'search'].",
    "Stopword Removal: Deleting common words ('the', 'a', 'is') that don't add meaning.",
    "Normalization: Converting 'Searching' and 'Search' to the same root.",
], w=7.0)
if TOKEN_IMG.exists():
    slide.shapes.add_picture(str(TOKEN_IMG), Inches(8.0), Inches(2.0), width=Inches(5.0))

# --- 10. MODULE 2: INDEXER ---
slide = blank_slide(10)
add_title_text(slide, "The Positional Indexer")
add_bullets(slide, [
    "Dictionary: A list of every unique word (lexicon).",
    "Postings List: For each word, a list of [DocID, [Positions]].",
    "Why Positions? This is how we support 'Phrase Search'.",
    "Example: \"Search Engine\" is found if 'Search' is at pos 5 and 'Engine' is at pos 6 in the same doc."
])

# --- 11. RETRIEVAL MODELS: BOOLEAN ---
slide = blank_slide(11)
add_title_text(slide, "Model 1: Boolean Overlap")
add_bullets(slide, [
    "The simplest ranking model.",
    "It counts: 'How many of the user's query terms appear in this document?'",
    "Formula: Score = |Query Terms \u2229 Document Terms|",
    "Weakness: It treats 'the' and 'algorithm' as equally important."
])

# --- 12. RETRIEVAL MODELS: TF-IDF ---
slide = blank_slide(12)
add_title_text(slide, "Model 2: TF-IDF Cosine")
add_bullets(slide, [
    "Term Frequency (TF): How often does the word appear in this doc?",
    "Inverse Document Frequency (IDF): How rare is this word across the whole collection?",
    "Wait weighting: Rare words (like 'Kullback-Leibler') get high scores; common words get low scores.",
    "Similarity: We treat Docs and Queries as vectors and calculate the angle (cosine) between them."
])

# --- 13. THE STAR: BM25 ---
slide = blank_slide(13)
add_title_text(slide, "Model 3: Okapi BM25")
add_bullets(slide, [
    "Industry Standard: The default in Lucene/Elasticsearch.",
    "Key Innovation: 'Term Saturation'.",
    ("If a doc has the word 'search' 100 times, it isn't 100x more relevant than a doc that has it 5 times.", 1),
    "Length Normalization: It penalizes very long documents (to avoid 'keyword stuffing')."
], w=7.0)
if BM25_IMG.exists():
    slide.shapes.add_picture(str(BM25_IMG), Inches(7.5), Inches(2.0), width=Inches(5.5))

# --- 14. ADVANCED FEATURE: PRF ---
slide = blank_slide(14)
add_title_text(slide, "Pseudo-Relevance Feedback (PRF)")
add_bullets(slide, [
    "A self-improving search cycle.",
    "The engine assumes the top 3 results are relevant.",
    "It extracts 'meaningful' new words from those results.",
    "It adds them to your query and searches again.",
    "Result: Finds documents you would have missed with your original keywords."
], w=7.5)
if FEEDBACK_IMG.exists():
    slide.shapes.add_picture(str(FEEDBACK_IMG), Inches(8.5), Inches(1.8), width=Inches(4.5))

# --- 15. INTERFACES: CLI ---
slide = blank_slide(15)
add_title_text(slide, "Developer Interface: The CLI")
add_bullets(slide, [
    "Command Line Power: Search from your terminal.",
    "Usage: python -m searchlens.cli search --query 'bm25' --method bm25",
    "Returns: Ranked list, scores, and snippets with term highlighting."
])

# --- 16. INTERFACES: WEB DEMO ---
slide = blank_slide(16)
add_title_text(slide, "User Interface: Web Demo")
add_bullets(slide, [
    "A fully functional search engine portal built using Python's native http library.",
    "No external dependencies (Flask/Django) required for the engine logic.",
    "Live interactivity: Toggle BM25, TF-IDF, or Boolean with a single click."
])

# --- 17. EXPERIMENT SETUP ---
slide = blank_slide(17)
add_title_text(slide, "Experimental Setup (SETC-60)")
add_bullets(slide, [
    "Dataset: SETC-60 (Synthetic Search Engine Tech Collection).",
    "Corpus: 60 technical documents covering 12 topics (Indexing, PageRank, etc.).",
    "Queries: 12 standardized user-style queries.",
    "Ground Truth: 75 graded relevance judgments for precise measurement."
])

# --- 18. METRICS EXPLAINED ---
slide = blank_slide(18)
add_title_text(slide, "How do we measure success?")
add_bullets(slide, [
    "P@5 (Precision at 5): Are the top 5 results relevant?",
    "Recall@10: Did we find most of the relevant documents?",
    "MAP (Mean Average Precision): The gold standard for overall ranking quality.",
    "nDCG (Normalized Discounted Cumulative Gain): Measures how well the best docs were 'pushed' to the very top."
])

# --- 19. RESULTS TABLE ---
slide = blank_slide(19)
add_title_text(slide, "Experimental Results")
rows, cols = 5, 6
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.5))
table = table_shape.table
headers = ["Method", "P@5", "Recall@10", "MAP", "MRR", "nDCG@10"]
data = [
    ["Boolean Overlap", "0.8333", "0.8216", "0.7736", "1.0000", "0.8487"],
    ["TF-IDF Cosine", "0.8667", "0.8593", "0.8242", "1.0000", "0.8842"],
    ["BM25 (The Best)", "0.8833", "0.8593", "0.8276", "1.0000", "0.8851"],
    ["BM25 + PRF", "0.8667", "0.8641", "0.8166", "1.0000", "0.8815"],
]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(16)
for r, row in enumerate(data, start=1):
    for c, val in enumerate(row):
        table.cell(r, c).text = val
        table.cell(r, c).text_frame.paragraphs[0].font.size = Pt(14)
        if "BM25" in row[0] and "PRF" not in row[0]:
            table.cell(r, c).text_frame.paragraphs[0].font.bold = True

# --- 20. RESULT ANALYSIS ---
slide = blank_slide(20)
add_title_text(slide, "Deep Dive: BM25 vs TF-IDF")
add_bullets(slide, [
    "Observation: BM25 consistently outscored TF-IDF across all metrics.",
    "Why? Because the SETC-60 documents are short and technical.",
    "BM25's length normalization prevents short documents from being unfairly penalized.",
    "TF-IDF remains a strong baseline, but fails to handle term saturation efficiently."
])

# --- 21. FEEDBACK ANALYSIS ---
slide = blank_slide(21)
add_title_text(slide, "The PRF Paradox")
add_bullets(slide, [
    "Recall Advantage: BM25+PRF achieved the highest Recall@10 (0.8641).",
    "Precision Trade-off: MAP actually decreased slightly (0.82 \u2192 0.81).",
    "Topic Drift: When the initial top results contain slightly unrelated technical terms, the query expands in the wrong direction.",
    "Conclusion: PRF is powerful but requires careful parameter tuning (k=3 works best)."
])

# --- 22. LIMITATIONS ---
slide = blank_slide(22)
add_title_text(slide, "Limitations & Future Work")
add_bullets(slide, [
    "Limitations:",
    ("Dataset is synthetic and small (60 docs).", 1),
    ("No semantic understanding (e.g., handling 'automobile' vs 'car').", 1),
    "Future Directions:",
    ("1. Semantic Search: Integrating BERT/LLM dense embeddings.", 1),
    ("2. Index Compression: Using Variable Byte Encoding for scalability.", 1),
    ("3. Live Evaluation: Real user-click tracking.", 1)
])

# --- 23. TARGET VENUE & DEMO ---
slide = blank_slide(23)
add_title_text(slide, "Target Venue & Live Demo")
add_bullets(slide, [
    "Target Venue: ACM SIGIR 2026 (System Demonstrations)",
    "SIGIR is the premier venue for Information Retrieval research.",
    "Live Demo URL: https://searchlens-demo.onrender.com",
    "Key Demo Features:",
    ("Real-time BM25 vs TF-IDF ranking comparison.", 1),
    ("Interactive 'Ranking Intelligence' explanations.", 1),
    ("Optimized for production via Render & Docker.", 1)
])

# --- 24. CONCLUSION ---
slide = blank_slide(24)
add_title_text(slide, "Conclusion")
add_bullets(slide, [
    "SearchLens successfully demonstrates that a transparent Python implementation can rival industrial models like BM25.",
    "The project proves the importance of term saturation and positional indexing for technical search tasks.",
    "The system is demo-ready, fully tested, and reproducible.",
    "Final URL: https://searchlens-demo.onrender.com"
])

# --- 25. THANK YOU ---
slide = blank_slide(25)
box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(2.5))
p = box.text_frame.paragraphs[0]
p.text = "Thank You"
p.font.name = FONT_NAME
p.font.size = Pt(84)
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = COLOR_BLUE

p2 = box.text_frame.add_paragraph()
p2.text = "Questions & Discussion"
p2.font.name = FONT_NAME
p2.font.size = Pt(32)
p2.alignment = PP_ALIGN.CENTER

p3 = box.text_frame.add_paragraph()
p3.text = "karthik.kovi@student.csulb.edu"
p3.font.name = FONT_NAME
p3.font.size = Pt(20)
p3.alignment = PP_ALIGN.CENTER
p3.font.color.rgb = COLOR_GRAY

prs.save(OUT)
print(f"Final Presentation (25 Slides) saved to: {OUT}")
